"""
Generate SurGen MIL presentation as a PowerPoint file.
Run from the repo root:  python scripts/make_presentation.py
Output: surgen_mil_presentation.pptx
"""

from pathlib import Path
import os
os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib")
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import textwrap

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
TEAL   = RGBColor(0x00, 0x7A, 0x87)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x55, 0x55, 0x55)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x00, 0x96, 0x88)  # teal accent
RED    = RGBColor(0xC6, 0x28, 0x28)

REPO = Path(__file__).parent.parent
FIG  = REPO / "docs" / "figures"
OUT  = REPO / "outputs" / "analysis" / "fair_comparison"
MULTI= REPO / "outputs" / "multisplit" / "analysis"
ATTN = REPO / "outputs" / "attention_viz"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = float(prs.slide_width)
H = float(prs.slide_height)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ──────────────────────────────────────────────────────────

def slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def bg(sl, color=WHITE):
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(sl, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txbox(sl, text, l, t, w, h,
          font_size=18, bold=False, color=DGRAY,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def multiline_txbox(sl, lines, l, t, w, h,
                    font_size=16, color=DGRAY, bold_first=False,
                    line_spacing=1.2, align=PP_ALIGN.LEFT):
    """lines: list of (text, bold, color_override) or just str."""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, (bold_first and first), color
        else:
            text, bold, col = item[0], item[1], (item[2] if len(item) > 2 else color)
        if first:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col
        first = False
    return tb

def header_bar(sl, title, subtitle=None, bar_h=1.15):
    rect(sl, 0, 0, 13.33, bar_h, fill=NAVY)
    txbox(sl, title, 0.35, 0.08, 12.0, 0.7,
          font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(sl, subtitle, 0.35, 0.72, 12.0, 0.38,
              font_size=16, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

def section_divider(sl, section_num, section_title, section_subtitle=""):
    rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    txbox(sl, f"Section {section_num}", 1.0, 2.2, 11.33, 0.5,
          font_size=20, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    txbox(sl, section_title, 1.0, 2.75, 11.33, 1.0,
          font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if section_subtitle:
        txbox(sl, section_subtitle, 1.0, 3.85, 11.33, 0.6,
              font_size=20, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

def img(sl, path, l, t, w, h=None):
    path = Path(path)
    if not path.exists():
        print(f"  [MISSING] {path}")
        return None
    if h is None:
        return sl.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
    else:
        return sl.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))

def bullet_list(sl, items, l, t, w, h, font_size=17, color=DGRAY,
                indent_char="▸ ", bold_labels=True):
    """items: list of str or (label, detail) tuples."""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(item, tuple):
            label, detail = item
            run = p.add_run()
            run.text = indent_char + label
            run.font.size = Pt(font_size)
            run.font.bold = bold_labels
            run.font.color.rgb = color
            run2 = p.add_run()
            run2.text = " " + detail
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = GRAY
        else:
            run = p.add_run()
            run.text = indent_char + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return tb

def caption(sl, text, l, t, w, h=0.35):
    txbox(sl, text, l, t, w, h, font_size=12, color=GRAY,
          italic=True, align=PP_ALIGN.CENTER)

def divider_line(sl, y, l=0.35, w=12.63):
    line = sl.shapes.add_connector(1,
        Inches(l), Inches(y), Inches(l+w), Inches(y))
    line.line.color.rgb = TEAL
    line.line.width = Pt(1.5)

def make_cohort_grouped_plot(out_path: Path):
    raw = MULTI / "raw_cohort_runs.csv"
    if not raw.exists():
        return None
    df = pd.read_csv(raw)
    keep = [
        "uni_mean_fair",
        "uni_attention_fair",
        "paper_reproduction_fair",
        "uni_hybrid_attention_mean2",
    ]
    label_map = {
        "uni_mean_fair": "MeanPool",
        "uni_attention_fair": "AttentionMIL",
        "paper_reproduction_fair": "TransformerMIL",
        "uni_hybrid_attention_mean2": "HybridAttentionMIL",
    }
    sub = df[df["model"].isin(keep)].copy()
    if sub.empty:
        return None
    agg = (
        sub.groupby(["model", "cohort"], as_index=False)
        .agg(mean=("cohort_auroc", "mean"), std=("cohort_auroc", "std"))
    )
    order = keep
    cohorts = ["SR1482", "SR386"]
    colors = {"SR1482": "#1f77b4", "SR386": "#ff7f0e"}
    x = range(len(order))
    width = 0.34

    fig, ax = plt.subplots(figsize=(9, 4.6))
    for j, cohort in enumerate(cohorts):
        vals = []
        errs = []
        for model in order:
            row = agg[(agg["model"] == model) & (agg["cohort"] == cohort)]
            vals.append(float(row["mean"].iloc[0]))
            errs.append(float(row["std"].iloc[0]) if not pd.isna(row["std"].iloc[0]) else 0.0)
        xpos = [i + (j - 0.5) * width for i in x]
        ax.bar(xpos, vals, width=width, color=colors[cohort], alpha=0.8, label=cohort)
        ax.errorbar(xpos, vals, yerr=errs, fmt="none", ecolor="black", elinewidth=1, capsize=3)

    split_df = pd.read_csv(MULTI / "summary_extended_by_split.csv")
    split_sub = split_df[split_df["model"].isin(keep)].copy()
    metric_col = {"SR1482": "SR1482_auroc", "SR386": "SR386_auroc"}
    for i, model in enumerate(order):
        for j, cohort in enumerate(cohorts):
            vals = split_sub.loc[split_sub["model"] == model, metric_col[cohort]].tolist()
            xpos = [i + (j - 0.5) * width] * len(vals)
            jitter = [-0.03, 0.0, 0.03][: len(vals)]
            ax.scatter([a + b for a, b in zip(xpos, jitter)], vals, color="black", s=18, zorder=3, alpha=0.8)

    ax.set_xticks(list(x))
    ax.set_xticklabels([label_map[m] for m in order], rotation=12, ha="right")
    ax.set_ylabel("Cohort AUROC")
    ax.set_title("Cohort AUROC by Model\nBars: mean across 9 runs, dots: split means")
    ax.set_ylim(0.75, 0.98)
    ax.grid(True, axis="y", alpha=0.25)
    ax.legend(frameon=False, ncol=2, loc="upper left")
    fig.tight_layout()
    fig.savefig(out_path, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return out_path

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
bg(sl, NAVY)
rect(sl, 0, 5.8, 13.33, 1.7, fill=TEAL)

txbox(sl, "SurGen MIL", 0.75, 0.9, 8.5, 0.82,
      font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txbox(sl, "Multiple Instance Learning for MSI/MMR Prediction", 0.75, 1.76, 8.5, 0.53,
      font_size=26, bold=False, color=LIGHT, align=PP_ALIGN.LEFT)
txbox(sl, "from Precomputed UNI Patch Embeddings · SurGen Colorectal Cancer Dataset",
      0.75, 2.29, 8.5, 0.41,
      font_size=18, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

divider_line(sl, 2.85, l=2.25, w=5.5)

txbox(sl, "Matthew Getzin", 0.75, 3.0, 8.5, 0.31,
      font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txbox(sl, "Computational Pathology · Take-Home Project for Pathomiq · March 2026",
      0.75, 3.38, 8.5, 0.26,
      font_size=15, color=LIGHT, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
bg(sl)
header_bar(sl, "Talk Outline")

sections = [
    ("1", "The Clinical Problem & Dataset",
         "MSI/MMR prediction in colorectal cancer from H&E slides"),
    ("2", "Starting Simple — Aggregation Strategies",
         "Why focus on aggregation when UNI embeddings already separate classes?"),
    ("3", "Data Challenges — Leakage & Label Normalization",
         "Case-grouped splits, SR1482 label reconciliation, cohort effects"),
    ("4", "The Fair Comparison",
         "MeanPool vs AttentionMIL vs TransformerMIL (paper reproduction)"),
    ("5", "Ablation Studies",
         "Loss functions, sparse evidence, sampler strategies"),
    ("6", "Attention Visualization",
         "Where are models looking? TP / FP / FN / TN maps"),
    ("7", "Multi-Split Stability & Extended Models",
         "9-model comparison · hybrid attention wins · spatial coordinate effects"),
    ("8", "Summary & Future Directions",
         "Key findings and what's next"),
]

for i, (num, title, sub) in enumerate(sections):
    row_t = 1.3 + i * 0.73
    rect(sl, 0.35, row_t, 0.55, 0.52, fill=NAVY)
    txbox(sl, num, 0.35, row_t, 0.55, 0.52,
          font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, title, 1.05, row_t, 5.2, 0.32,
          font_size=16, bold=True, color=NAVY)
    txbox(sl, sub, 1.05, row_t + 0.28, 12.0, 0.28,
          font_size=12, color=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 1 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 1, "The Clinical Problem & Dataset",
                "MSI/MMR status prediction in colorectal cancer")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Clinical context
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Why Predict MSI/MMR Status?",
           "Microsatellite instability / Mismatch repair deficiency in colorectal cancer")

bullet_list(sl, [
    ("MSI-H / dMMR",
     "~15% of CRC cases; defines immunogenic tumours that respond to checkpoint inhibitors (PD-1/PD-L1)."),
    ("Clinical gold standard",
     "IHC for MMR proteins + PCR-based MSI testing — costly, requires dedicated lab infrastructure."),
    ("H&E as a low-cost proxy",
     "Deep learning models can read morphological clues (TIL density, mucinous features, Crohn-like reactions) "
     "directly from routine H&E stained slides."),
    ("Weak supervision challenge",
     "Only slide-level labels available — patch-level ground truth absent. "
     "MIL frameworks are the natural fit."),
    ("Dataset: SurGen",
     "Two cohorts (SR1482 primary tumours, SR386 primary+mets), precomputed UNI patch embeddings "
     "in Zarr format — no WSI preprocessing required."),
], l=0.35, t=1.3, w=12.6, h=5.8, font_size=17)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Dataset
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Dataset: SurGen Colorectal Cancer Cohorts")

# Left column — SR1482
rect(sl, 0.35, 1.3, 5.9, 5.6, fill=RGBColor(0xE8, 0xF4, 0xF8),
     line_color=TEAL, line_width=Pt(1))
txbox(sl, "SR1482 — Primary Tumours", 0.5, 1.45, 5.6, 0.45,
      font_size=17, bold=True, color=NAVY)
bullet_list(sl, [
    "Primary colorectal carcinoma slides",
    "Labels: MSI column + MMR IHC columns — reconciled via 8-state logic",
    "Discordant / inconclusive cases excluded",
    "Each case may have multiple slides (same case → same split)",
    "1024-dim UNI embeddings, 40× magnification",
], l=0.5, t=1.95, w=5.6, h=4.5, font_size=15, indent_char="• ")

# Right column — SR386
rect(sl, 6.6, 1.3, 6.38, 5.6, fill=RGBColor(0xE8, 0xF8, 0xEE),
     line_color=GREEN, line_width=Pt(1))
txbox(sl, "SR386 — Primary + Metastases", 6.75, 1.45, 6.1, 0.45,
      font_size=17, bold=True, color=GREEN)
bullet_list(sl, [
    "Includes primary and metastatic CRC slides",
    "Label: mmr_loss_binary column — direct binary target",
    "Simpler label structure but different class prevalence",
    "Same UNI embedding format",
    "Cohort effects present — modelled jointly but monitored separately",
], l=6.75, t=1.95, w=6.1, h=4.5, font_size=15, indent_char="• ")

# Centre label
txbox(sl, "Combined cohort · Case-grouped stratified train / val / test split (split_seed=0)",
      0.35, 6.7, 12.6, 0.35, font_size=13, color=GRAY, italic=True,
      align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 2 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 2, "Starting Simple — Aggregation Strategies",
                "UNI embeddings already do heavy lifting")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Design philosophy
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "First Design Decision: Focus on Aggregation",
           "Feature extraction is already handled by frozen UNI — what's the best way to aggregate?")

txbox(sl, "UNI (Universal Pathology Foundation Model) provides precomputed 1024-dim patch embeddings. "
      "The upstream representation problem is solved. The question becomes: "
      "how do we best aggregate hundreds to thousands of patch embeddings into a single slide-level prediction?",
      0.35, 1.35, 12.6, 1.0, font_size=17, color=DGRAY)

divider_line(sl, 2.45)

# Two columns
txbox(sl, "Why not start with a more complex model?", 0.35, 2.6, 6.0, 0.4,
      font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "Limited labelled cases — weak supervision risk of overfitting",
    "Frozen embeddings remove the need for backprop through a backbone",
    "Aggregation is the only trainable step",
    "Start simple, quantify gain from complexity",
], l=0.35, t=3.05, w=6.0, h=3.5, font_size=16, indent_char="▸ ")

txbox(sl, "Aggregation hypothesis", 6.7, 2.6, 6.0, 0.4,
      font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "Mean pooling: if embeddings are linearly separable, a linear probe suffices",
    "Attention MIL: learns which patches are diagnostically relevant",
    "Transformer: models patch-to-patch interactions (more expressive, higher risk)",
    "Paper baseline: 2-layer transformer encoder (~6.8M params)",
], l=6.7, t=3.05, w=6.0, h=3.5, font_size=16, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Pipeline overview
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "MIL Pipeline",
           "Modular design: swap any component independently")

# Pipeline boxes
steps = [
    ("UNI Patch\nEmbeddings", "N × 1024\n(precomputed)", TEAL),
    ("Bag\nSampler", "Random / spatial /\nfeature-diverse", NAVY),
    ("Aggregator", "MeanPool /\nAttentionMIL /\nTransformerMIL", ACCENT),
    ("Classifier\nHead", "MLP → logit", NAVY),
    ("Loss +\nOptimizer", "BCE (weighted) /\nFocal · AdamW", TEAL),
]
box_w = 2.1
gap = 0.2
start_l = 0.35
for i, (label, detail, col) in enumerate(steps):
    l = start_l + i * (box_w + gap)
    rect(sl, l, 2.2, box_w, 1.8, fill=col)
    txbox(sl, label, l, 2.35, box_w, 0.7,
          font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, detail, l, 2.9, box_w, 1.0,
          font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        # Arrow
        ax = l + box_w + 0.02
        txbox(sl, "→", ax, 2.85, gap + 0.15, 0.4,
              font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txbox(sl, "Evaluation", 0.35, 4.35, 12.6, 0.38,
      font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    ("AUROC / AUPRC", "slide-level · temperature-scaled predictions · EMA checkpoint selection on val AUPRC"),
    ("Confusion matrix", "with Youden-J threshold from validation set"),
    ("Case-level aggregation", "max / mean / noisy-or across slides from same patient"),
], l=0.35, t=4.75, w=12.6, h=2.3, font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Three architectures
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Three Aggregation Architectures",
           "Designed for a fair head-to-head comparison")

cols = [
    ("MeanPool", "262K params",
     ["Project: 1024 → 256 (ReLU, dropout)",
      "Mean pool across all patches",
      "Classifier: 256 → 1 (linear)",
      "",
      "Hypothesis: frozen embeddings are",
      "linearly separable — mean is enough.",
      "Lowest variance, fastest training."],
     TEAL),
    ("AttentionMIL", "394K params",
     ["Attention scorer: 1024 → 128 → 1",
      "Softmax over patches → weights",
      "Weighted sum → 1024-dim context",
      "Classifier: 1024 → 1 (linear)",
      "",
      "Selectively up-weights diagnostic patches.",
      "Higher capacity, higher variance."],
     NAVY),
    ("TransformerMIL\n(Paper Repro)", "6.8M params",
     ["Project: 1024 → 512",
      "2-layer transformer encoder",
      "(4 heads, 2048 FFN)",
      "Mean pool → Classifier: 512 → 1",
      "",
      "Matches paper architecture.",
      "Most expressive — risk of overfitting."],
     RGBColor(0x4A, 0x14, 0x8C)),
]
col_w = 3.9
for i, (name, params, desc, col) in enumerate(cols):
    l = 0.35 + i * (col_w + 0.3)
    rect(sl, l, 1.25, col_w, 0.65, fill=col)
    txbox(sl, name, l, 1.28, col_w, 0.42,
          font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, params, l, 1.72, col_w, 0.3,
          font_size=13, color=col, align=PP_ALIGN.CENTER, bold=True)
    for j, line in enumerate(desc):
        txbox(sl, ("▸ " if line and not line.startswith("Hypo") else "") + line,
              l + 0.1, 2.1 + j * 0.56, col_w - 0.2, 0.5,
              font_size=14, color=DGRAY if not line.startswith("Hypo") and line else GRAY,
              italic=line.startswith("Hypo") or not line)


# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 3 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 3, "Data Challenges",
                "Leakage, label noise, and cohort normalization")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Leakage
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Data Leakage: The Case-Grouping Problem",
           "A subtle but critical issue in multi-slide datasets")

txbox(sl, "The Problem", 0.35, 1.3, 12.6, 0.38,
      font_size=17, bold=True, color=RED)
txbox(sl, "Multiple slides from the same patient (case) can end up in different splits "
      "under naive random stratification. This leaks patient-level information across "
      "train / val / test boundaries — artificially inflating val/test metrics.",
      0.35, 1.7, 12.6, 0.75, font_size=16, color=DGRAY)

divider_line(sl, 2.6)

txbox(sl, "The Fix", 0.35, 2.7, 12.6, 0.38,
      font_size=17, bold=True, color=GREEN)
bullet_list(sl, [
    ("Case-grouped split",
     "All slides from a given case_id are assigned to the same split before stratification."),
    ("Stratification within cases",
     "Splits are balanced by label distribution across cases (not slides), "
     "ensuring class prevalence is comparable across train/val/test."),
    ("Naming convention",
     "Slide IDs encode case_id: {COHORT}_40X_HE_T{CASE_ID}_{SLIDE_NUM} — "
     "grouping is automatic."),
    ("Impact",
     "Discovered during early exploratory runs; corrected before the main fair-comparison experiments. "
     "Earlier (leaky) results showed inflated AUROC — discarded."),
], l=0.35, t=3.1, w=12.6, h=3.8, font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Label normalization
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Cohort Label Normalization",
           "SR1482 requires reconciliation across MSI and MMR annotation sources")

txbox(sl, "SR1482 — Complex 8-State Reconciliation", 0.35, 1.3, 6.3, 0.4,
      font_size=16, bold=True, color=NAVY)

states = [
    ("MSI-H + dMMR", "→ Include as positive (1)", GREEN),
    ("MSS + pMMR", "→ Include as negative (0)", TEAL),
    ("MSI-H + pMMR", "→ Discordant — EXCLUDE", RED),
    ("MSS + dMMR", "→ Discordant — EXCLUDE", RED),
    ("MSI column missing", "→ Fall back to MMR only", GRAY),
    ("MMR column missing", "→ Fall back to MSI only", GRAY),
    ("Both missing", "→ EXCLUDE (no label)", RED),
    ("Inconclusive", "→ EXCLUDE", RED),
]
for i, (state, action, col) in enumerate(states):
    row = 1.78 + i * 0.52
    txbox(sl, state, 0.35, row, 3.8, 0.45, font_size=14, bold=True, color=col)
    txbox(sl, action, 4.2, row, 2.3, 0.45, font_size=14, color=DGRAY)

txbox(sl, "SR386 — Simple Binary", 6.85, 1.3, 6.0, 0.4,
      font_size=16, bold=True, color=GREEN)
txbox(sl, "mmr_loss_binary column provides a direct 0/1 label. "
      "No reconciliation needed, but class balance differs from SR1482.",
      6.85, 1.78, 6.0, 1.0, font_size=15, color=DGRAY)

txbox(sl, "Why This Matters", 6.85, 2.95, 6.0, 0.4,
      font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "Discordant IHC/PCR cases introduce ambiguous training signal",
    "Excluding them reduces dataset size but improves label quality",
    "Per-cohort class balance differs → weighted loss critical",
    "Monitored cohort-stratified metrics throughout",
], l=6.85, t=3.4, w=6.0, h=3.5, font_size=15, indent_char="▸ ")


# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 4 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 4, "The Fair Comparison",
                "uni_mean_fair · uni_attention_fair · paper_reproduction_fair")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Experimental setup
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Fair Comparison Setup",
           "After fixing leakage and label normalization — an apples-to-apples test")

bullet_list(sl, [
    ("Fixed data split", "split_seed=0 · case-grouped stratified · same split for all three models"),
    ("Three training seeds", "42, 123, 456 — captures run-to-run variance"),
    ("Weighted BCE loss", "class weights from training set label frequency"),
    ("AdamW optimizer", "lr=1e-4, weight_decay=1e-4 · 50 epochs"),
    ("Checkpoint selection", "EMA-smoothed val AUPRC (α=0.7) — reduces noise from single-epoch evaluation"),
    ("Temperature scaling", "post-hoc calibration on validation set · minimises NLL"),
    ("Evaluation", "slide-level AUROC and AUPRC on held-out test set · case-level aggregation (max/mean/noisy_or)"),
    ("Paper baseline", "AUROC 0.827, but not directly comparable: different split, 200 epochs, "
     "no early stopping / EMA checkpointing, and no random train-time bag sampling"),
], l=0.35, t=1.3, w=12.6, h=5.8, font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Main results table
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Main Results: Slide-Level Performance",
           "Mean ± std across 3 seeds · All models exceed paper baseline AUROC 0.827")

# Table
headers = ["Model", "Params", "AUROC (mean ± std)", "AUPRC (mean ± std)", "Interpretation"]
rows = [
    ["MeanPool", "262K", "0.860 ± 0.005", "0.406 ± 0.019",
     "Most stable · lowest variance · strong linear probe"],
    ["AttentionMIL", "394K", "0.869 ± 0.020", "0.420 ± 0.032",
     "Marginally higher mean · 4× variance · seed-dependent"],
    ["TransformerMIL\n(paper repro)", "6.8M", "0.806 ± 0.057", "0.349 ± 0.072",
     "Lowest AUROC · highest variance · likely overfitting"],
    ["Paper baseline", "6.8M", "0.827 (reported)", "—",
     "Different split — reference only"],
]
col_ws = [2.0, 0.9, 2.2, 2.2, 4.65]
col_ls = [0.35]
for w in col_ws[:-1]:
    col_ls.append(col_ls[-1] + w + 0.05)

row_h = 0.75
header_t = 1.35
for j, (h, l, w) in enumerate(zip(headers, col_ls, col_ws)):
    rect(sl, l, header_t, w, 0.5, fill=NAVY)
    txbox(sl, h, l + 0.05, header_t + 0.05, w - 0.1, 0.45,
          font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

row_colors = [RGBColor(0xE8, 0xF8, 0xEE), RGBColor(0xE8, 0xF4, 0xF8),
              RGBColor(0xFD, 0xE8, 0xE8), RGBColor(0xF5, 0xF5, 0xF5)]
for i, (row, rc) in enumerate(zip(rows, row_colors)):
    t = header_t + 0.5 + i * row_h
    for j, (cell, l, w) in enumerate(zip(row, col_ls, col_ws)):
        rect(sl, l, t, w, row_h - 0.05,
             fill=rc, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        is_metric = j in (2, 3) and i < 3
        txbox(sl, cell, l + 0.05, t + 0.05, w - 0.1, row_h - 0.15,
              font_size=13 if j != 0 else 13,
              bold=is_metric, color=GREEN if is_metric and i == 1 else
              (RED if is_metric and i == 2 else DGRAY),
              align=PP_ALIGN.CENTER if j in (1, 2, 3) else PP_ALIGN.LEFT)

txbox(sl, "★ Key finding: frozen UNI embeddings are already highly discriminative. "
      "Simple mean pooling is a strong baseline. Expressive architectures do not consistently win "
      "at this dataset scale under weak supervision.",
      0.35, 6.95, 12.6, 0.48, font_size=14, bold=True, color=NAVY,
      align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — ROC & PR curves
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "ROC and Precision-Recall Curves",
           "Fair comparison · 3 seeds per model · fixed split")

roc_path = OUT / "roc_pr.png"
if roc_path.exists():
    img(sl, roc_path, 0.35, 1.25, 12.6)
    caption(sl, "Left: ROC curves. Right: PR curves. Curves are seed-averaged test probabilities "
            "per model on the fixed fair-comparison split.",
            0.35, 6.82, 12.6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Confusion matrices
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Confusion Matrices — Fair Comparison",
           "Seed-averaged test probabilities · Youden J threshold from averaged validation predictions")

cm_path = OUT / "confusion_matrices.png"
if cm_path.exists():
    img(sl, cm_path, 0.35, 1.25, 8.4)
    caption(sl, "Top row: slide-level. Bottom row: case-level mean aggregation. "
            "Each panel uses the threshold fit on averaged validation probabilities for that model.",
            0.35, 6.82, 8.4)

txbox(sl, "What is averaged?", 9.0, 1.35, 3.7, 0.35,
      font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "Test probabilities are averaged across the 3 seeds for each slide",
    "Validation probabilities are averaged across the 3 seeds before fitting Youden J",
    "So the threshold is not a per-seed threshold; it is fit on the averaged validation prediction",
    "Positive-class scarcity is still visible: most operating points trade recall against a large negative pool",
], l=9.0, t=1.78, w=3.7, h=4.9, font_size=13, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Training curves + metric summary
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Training Dynamics & Metric Summary",
           "EMA-smoothed val AUPRC checkpoint selection · 50 epochs")

tc_path = OUT / "training_curves.png"
ms_path = OUT / "metric_summary.png"
if tc_path.exists():
    img(sl, tc_path, 0.35, 1.25, 7.15)
    caption(sl, "Training curves: rows = models, columns = val AUROC / val AUPRC / train loss",
            0.35, 6.82, 7.15)
if ms_path.exists():
    img(sl, ms_path, 7.85, 1.25, 5.15)
    caption(sl, "Metric summary across seeds", 7.85, 3.0, 5.15)

txbox(sl, "Takeaways", 7.85, 3.35, 5.15, 0.35,
      font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "Training curves: MeanPool is the tightest across seeds; AttentionMIL is noisier; TransformerMIL is the least stable",
    "Bar summary: MeanPool wins on stability, AttentionMIL edges the mean AUROC, and TransformerMIL loses on both mean and variance",
    "This is the cleanest place to discuss variance because both the per-seed curves and summary bars are visible together",
], l=7.85, t=3.75, w=5.15, h=2.7, font_size=13, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Interpretation
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "What Is Bottlenecking Performance?",
           "The bottleneck is aggregation stability, not representation")

# Three findings boxes
findings = [
    ("Representation is not the bottleneck",
     "UNI embeddings encode rich patch-level features. "
     "Even a linear mean pool achieves AUROC 0.860. "
     "The signal is there — the question is how reliably we extract it.",
     TEAL),
    ("Aggregation quality matters, but more is not always better",
     "AttentionMIL marginally outperforms MeanPool in mean AUROC "
     "but with 4× higher variance. "
     "TransformerMIL — with 17× more parameters — performs worst. "
     "Weak supervision and dataset size penalize overparameterization.",
     NAVY),
    ("Attention adds value at a cost",
     "Attention MIL can identify diagnostically relevant patches. "
     "But its optimization landscape is sensitive to weight initialization and seed. "
     "This variance is the central challenge for clinical deployment.",
     ACCENT),
]
for i, (title, body, col) in enumerate(findings):
    t = 1.35 + i * 1.9
    rect(sl, 0.35, t, 0.08, 1.55, fill=col)
    txbox(sl, title, 0.6, t, 12.1, 0.42,
          font_size=17, bold=True, color=col)
    txbox(sl, body, 0.6, t + 0.43, 12.1, 1.1,
          font_size=15, color=DGRAY)

txbox(sl, "→ Next: ablations to probe each of these hypotheses systematically.",
      0.35, 6.9, 12.6, 0.45, font_size=14, bold=True, color=NAVY,
      align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 5 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 5, "Ablation Studies",
                "Appendices A – D: loss, sparse evidence, samplers")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Appendix A: Pooling variant & loss weighting
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Appendix A — Pooling Variant & Loss Weighting",
           "Does it matter how we pool, or whether we weight the loss? · 3 seeds · split_seed=0")

# Table
def ablation_table(sl, title, col_headers, rows, row_colors,
                   col_ws, t_top, note=None, highlight_row=None):
    """Draw a simple ablation results table."""
    txbox(sl, title, 0.35, t_top, 12.6, 0.38, font_size=15, bold=True, color=NAVY)
    col_ls = [0.35]
    for w in col_ws[:-1]:
        col_ls.append(col_ls[-1] + w + 0.05)
    hdr_t = t_top + 0.42
    for j, (h, l, w) in enumerate(zip(col_headers, col_ls, col_ws)):
        rect(sl, l, hdr_t, w, 0.42, fill=NAVY)
        txbox(sl, h, l+0.05, hdr_t+0.05, w-0.1, 0.38,
              font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (row, rc) in enumerate(zip(rows, row_colors)):
        rt = hdr_t + 0.42 + i * 0.58
        is_hl = (i == highlight_row)
        for j, (cell, l, w) in enumerate(zip(row, col_ls, col_ws)):
            rect(sl, l, rt, w, 0.54, fill=rc,
                 line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.5))
            txbox(sl, cell, l+0.06, rt+0.06, w-0.12, 0.44,
                  font_size=13, bold=(is_hl and j > 1),
                  color=(GREEN if is_hl and j > 1 else DGRAY),
                  align=PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT)
    if note:
        note_t = hdr_t + 0.42 + len(rows) * 0.58 + 0.08
        txbox(sl, note, 0.35, note_t, 12.6, 0.45,
              font_size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

ablation_table(
    sl,
    "Results (mean ± std, 3 seeds)",
    ["Model", "Aggregation", "Loss", "AUROC", "AUPRC"],
    [
        ["MeanPool",      "Bag-level mean",    "Weighted BCE",   "0.860 ± 0.005", "0.447 ± 0.019"],
        ["MeanPool",      "Bag-level mean",    "Unweighted BCE", "0.862 ± 0.003", "0.465 ± 0.004"],
        ["Instance Mean", "Classify-then-pool","Weighted BCE",   "0.859 ± 0.008", "0.443 ± 0.027"],
    ],
    [RGBColor(0xE8,0xF8,0xEE), RGBColor(0xF5,0xF5,0xF5), RGBColor(0xF0,0xF4,0xFB)],
    col_ws=[2.5, 2.8, 2.6, 2.05, 2.05],
    t_top=1.3,
    note="All three variants are within noise of each other. The dominant factor is the representation, not the pooling or loss weighting.",
)

divider_line(sl, 4.2)
txbox(sl, "Interpretation", 0.35, 4.3, 12.6, 0.38, font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    ("Unweighted BCE is not worse:",
     "AUROC 0.862 ± 0.003 vs 0.860 ± 0.005 — class weighting has minimal impact at this imbalance ratio. "
     "AUPRC is slightly higher, but this slide alone does not support a confusion-matrix claim about error type."),
    ("Instance mean matches bag mean:",
     "Classifying each patch independently then pooling probabilities (instance mean) "
     "yields the same result as pooling embeddings first — further evidence the embeddings are "
     "already well-calibrated in feature space."),
    ("Conclusion:",
     "Neither pooling strategy nor loss weighting moves the needle. Aggregation architecture is the variable to explore."),
], l=0.35, t=4.72, w=12.6, h=2.5, font_size=14)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — Appendix B: Focal loss
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Appendix B — Focal Loss vs Weighted BCE for AttentionMIL",
           "Use the figures to talk about score distributions and calibration, not just the metric table")

ablation_table(
    sl,
    "Results (mean ± std, 3 seeds)",
    ["Model", "Loss", "AUROC", "AUPRC", "AUPRC std"],
    [
        ["AttentionMIL", "Weighted BCE",  "0.869 ± 0.020", "0.381 ± 0.052", "±0.052"],
        ["AttentionMIL", "Focal (γ=2)",  "0.864 ± 0.014", "0.420 ± 0.066", "±0.066"],
    ],
    [RGBColor(0xE8,0xF4,0xF8), RGBColor(0xF0,0xF4,0xFB)],
    col_ws=[2.6, 2.5, 2.4, 2.4, 1.95],
    t_top=1.3,
    note="Focal loss trades ~0.5pp AUROC for ~4pp AUPRC gain. The plots below are better evidence for how the score behaviour changes.",
)

cal_path = OUT / "calibration.png"
err_path = OUT / "error_distributions.png"
if cal_path.exists():
    img(sl, cal_path, 0.35, 3.65, 4.15)
    caption(sl, "Calibration after temperature scaling", 0.35, 6.82, 4.15)
if err_path.exists():
    img(sl, err_path, 4.8, 3.65, 5.1)
    caption(sl, "Seed-averaged score distributions and thresholds", 4.8, 6.82, 5.1)

txbox(sl, "What the figures support", 10.15, 3.65, 2.8, 0.35, font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "Focal loss changes score shape, not just one summary metric",
    "AUPRC rises, but variance also rises: ±0.066 vs ±0.052",
    "Use this framing for a recall-oriented use case, not a blanket claim about calibration",
], l=10.15, t=4.05, w=2.8, h=2.7, font_size=12, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — Appendix C: Top-k sparse evidence
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Appendix C — Sparse Evidence: Top-k = 16",
           "Select the 16 highest-attended patches at test time — does focus help?")

# Top row: table + figure
ablation_table(
    sl,
    "Results (mean ± std, 3 seeds)",
    ["Variant", "Patches used", "AUROC", "AUPRC"],
    [
        ["AttentionMIL", "Full bag (~512)", "0.869 ± 0.020", "0.381 ± 0.052"],
        ["AttentionMIL (top-k=16)", "Top 16 only", "0.853 ± 0.032", "0.455 ± 0.140"],
    ],
    [RGBColor(0xE8,0xF4,0xF8), RGBColor(0xF5,0xF0,0xFD)],
    col_ws=[3.2, 2.6, 2.1, 2.1],
    t_top=1.3,
)

c_roc = FIG / "appendix_c_roc_pr_curves.png"
if c_roc.exists():
    img(sl, c_roc, 6.85, 1.3, 6.0)
    caption(sl, "ROC / PR curves · full-bag vs top-k=16", 6.85, 4.05, 6.0)

divider_line(sl, 4.35)

txbox(sl, "Interpretation", 0.35, 4.45, 12.6, 0.38, font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "AUROC drops −1.6pp: forcing focus on 16 patches discards contextual signal spread across the slide",
    "AUPRC gains +7.4pp: model becomes more precise on MSI-H — but AUPRC std explodes (±0.140 vs ±0.052)",
    "Variance increase suggests top-k is fragile: whether the right 16 patches are selected depends heavily on initialization",
    "Clinical use case: when a concise heatmap with few highlighted regions is needed for pathologist review",
], l=0.35, t=4.88, w=12.6, h=1.95, font_size=13, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19 — Appendix D: Sampler strategies
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Appendix D — Train-Time Sampler Strategies",
           "Does the patch sampling strategy during training affect the learned aggregator? · full-bag eval")

def mini_table(sl, x0, title, rows, row_colors):
    txbox(sl, title, x0, 1.3, 5.9, 0.32, font_size=15, bold=True,
          color=TEAL if "MeanPool" in title else NAVY)
    headers = ["Sampler", "AUROC", "AUPRC", "std"]
    widths = [2.2, 1.35, 1.35, 0.8]
    xs = [x0]
    for w in widths[:-1]:
        xs.append(xs[-1] + w + 0.05)
    y = 1.68
    for h, x, w in zip(headers, xs, widths):
        rect(sl, x, y, w, 0.36, fill=NAVY)
        txbox(sl, h, x + 0.04, y + 0.04, w - 0.08, 0.28,
              font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (row, rc) in enumerate(zip(rows, row_colors)):
        ry = y + 0.39 + i * 0.46
        for cell, x, w in zip(row, xs, widths):
            rect(sl, x, ry, w, 0.42, fill=rc,
                 line_color=RGBColor(0xCC,0xCC,0xCC), line_width=Pt(0.4))
            txbox(sl, cell, x + 0.04, ry + 0.05, w - 0.08, 0.3,
                  font_size=11, color=DGRAY,
                  align=PP_ALIGN.CENTER if x != xs[0] else PP_ALIGN.LEFT)

mini_table(
    sl, 0.35, "MeanPool (mean ± std, 3 seeds)",
    [
        ["Random",          "0.860 ± 0.005", "0.447 ± 0.019", "±0.019"],
        ["Spatial grid",    "0.859 ± 0.007", "0.447 ± 0.015", "±0.015"],
        ["Feature-diverse", "0.852 ± 0.005", "0.358 ± 0.035", "±0.035"],
    ],
    [RGBColor(0xEE,0xF8,0xEE), RGBColor(0xF5,0xF5,0xF5), RGBColor(0xF5,0xF5,0xF5)],
)

mini_table(
    sl, 6.85, "AttentionMIL (mean ± std, 3 seeds)",
    [
        ["Random",          "0.869 ± 0.020", "0.381 ± 0.052", "±0.052"],
        ["Spatial grid",    "0.861 ± 0.031", "0.404 ± 0.058", "±0.058"],
        ["Feature-diverse", "0.879 ± 0.006", "0.407 ± 0.058", "±0.058"],
    ],
    [RGBColor(0xE8,0xF4,0xF8), RGBColor(0xF5,0xF5,0xF5), RGBColor(0xD4,0xED,0xDA)],
)

divider_line(sl, 4.1)

txbox(sl, "MeanPool", 0.35, 4.2, 5.9, 0.35, font_size=15, bold=True, color=TEAL)
bullet_list(sl, [
    "AUROC is stable across all three samplers (within 0.008) — mean pool is permutation-invariant by design",
    "Feature-diverse drops AUPRC to 0.358: diverse sampling may under-represent the rare positive patches that drive recall",
], l=0.35, t=4.58, w=5.9, h=2.2, font_size=13, indent_char="▸ ")

txbox(sl, "AttentionMIL", 6.7, 4.2, 6.0, 0.35, font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "Feature-diverse: AUROC 0.879 ± 0.006 — best AUROC, and variance collapses from ±0.020 to ±0.006",
    "Spatial grid: intermediate — reduces variance somewhat but not as cleanly",
    "Explanation: diverse training bags force attention to generalize across feature types rather than latching onto an easy, repeating pattern",
], l=6.7, t=4.58, w=6.0, h=2.2, font_size=13, indent_char="▸ ")

# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 6 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 6, "Attention Visualization",
                "Representative success and failure cases from the core fair-comparison models")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE: Representative successes
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Attention Maps — Representative Successes",
           "Left to right in each panel: H&E, experimental attention, MeanPool IG, AttentionMIL, TransformerMIL IG")

tp_fig = FIG / "attn_tp_SR386_T402.png"
tn_fig = FIG / "attn_tn_SR386_T129.png"
if tp_fig.exists():
    img(sl, tp_fig, 0.35, 1.25, 12.6)
    caption(sl, "True positive: all models confidently identify MSI-H", 0.35, 3.55, 12.6)
if tn_fig.exists():
    img(sl, tn_fig, 0.35, 3.95, 12.6)
    caption(sl, "True negative: all models correctly suppress the positive class", 0.35, 6.32, 12.6)

txbox(sl, "Takeaways", 0.35, 6.6, 12.6, 0.3, font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "These examples only use models already introduced in the fair comparison; hybrid/spatial variants appear later in section 7",
    "When the case is easy, the maps are broadly consistent across methods: distributed positive evidence in the TP, uniformly low attribution in the TN",
    "MeanPool integrated gradients often track the learned attention maps, which argues the signal is not purely an attention artifact",
], l=0.35, t=6.92, w=12.6, h=0.45, font_size=12, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE: Representative failures
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Attention Maps — Representative Failures",
           "Same left-to-right order: H&E, experimental attention, MeanPool IG, AttentionMIL, TransformerMIL IG")

fp_fig = FIG / "attn_fp_SR1482_T061.png"
fn_fig = FIG / "attn_fn_SR386_T436.png"
if fp_fig.exists():
    img(sl, fp_fig, 0.35, 1.25, 12.6)
    caption(sl, "False positive: all models are confidently wrong on the same slide", 0.35, 3.55, 12.6)
if fn_fig.exists():
    img(sl, fn_fig, 0.35, 3.95, 12.6)
    caption(sl, "False negative: some attention signal is present, but the fair-comparison models do not recover it reliably", 0.35, 6.32, 12.6)

txbox(sl, "Takeaways", 0.35, 6.6, 12.6, 0.3, font_size=15, bold=True, color=NAVY)
bullet_list(sl, [
    "The false positive is a systematic morphology-mimic error, not a one-seed anomaly",
    "The false negative shows that signal can be present yet unstable under weak supervision; attention-based aggregation sometimes recovers it, but not reliably enough",
    "The failure story motivates section 7: improve stability first, then add architectural complexity",
], l=0.35, t=6.92, w=12.6, h=0.45, font_size=12, indent_char="▸ ")

# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 7 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 7, "Multi-Split Stability & Extended Models",
                "3 splits × 3 seeds · hybrid attention · spatial variants")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 21 — Multisplit design
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Multi-Split Experimental Design",
           "3 data splits × 3 seeds = 9 runs per model — robust performance estimates")

txbox(sl, "Motivation", 0.35, 1.35, 12.6, 0.38,
      font_size=16, bold=True, color=NAVY)
txbox(sl, "A single train/val/test split may be particularly favourable or unfavourable "
      "for a given model. Multi-split evaluation averages over split randomness, giving "
      "a more reliable estimate of expected performance on unseen data.",
      0.35, 1.78, 12.6, 0.75, font_size=16, color=DGRAY)

divider_line(sl, 2.65)

bullet_list(sl, [
    ("Splits:", "split_seed ∈ {0, 1, 2} — each generates a different case-grouped stratified partition"),
    ("Seeds:", "training seed ∈ {42, 123, 456} per split — same 3 seeds as fair comparison"),
    ("Core models:", "uni_mean_fair · uni_attention_fair · paper_reproduction_fair"),
    ("Extended models:", "uni_mean_var · uni_gated_attention · uni_attention_spatial_fair · "
     "uni_hybrid_attention_mean2 · uni_hybrid_attention_spatial_mean2"),
    ("Spatial TransformerMIL:", "New config uni_transformer_spatial_fair: CoordinateEncoder produces "
     "32-dim embedding concatenated to each patch (1024+32 → 512 projection). "
     "Gives the transformer explicit spatial grounding — evaluated last."),
    ("Aggregation:", "Mean AUROC / AUPRC across 9 runs per model; cohort-stratified strips computed separately"),
], l=0.35, t=2.82, w=12.6, h=4.35, font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 22 — Multisplit full results table
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Multi-Split Results — All Models",
           "9 runs per model (3 splits × 3 seeds) · slide-level AUROC and AUPRC")

# Table data — ranked by AUROC
ms_headers = ["Model", "AUROC (mean ± std)", "AUPRC (mean ± std)", "Notes"]
ms_rows = [
    ["HybridAttentionMIL\n(mean + 2 heads)", "0.903 ± 0.035", "0.591 ± 0.058",
     "★ Best overall — highest AUROC and AUPRC"],
    ["AttentionMIL", "0.900 ± 0.032", "0.532 ± 0.128",
     "Strong but high AUPRC variance (±0.128)"],
    ["Hybrid + Spatial", "0.897 ± 0.040", "0.541 ± 0.074",
     "Coords add little over non-spatial hybrid"],
    ["Gated AttentionMIL", "0.896 ± 0.040", "0.516 ± 0.138",
     "Gating mechanism — higher AUPRC variance"],
    ["MeanVar Pool", "0.895 ± 0.027", "0.497 ± 0.084",
     "Mean + variance — stable, no attention needed"],
    ["Attention + Spatial", "0.882 ± 0.041", "0.485 ± 0.146",
     "Coords hurt: underperforms plain attention"],
    ["MeanPool", "0.877 ± 0.033", "0.495 ± 0.050",
     "Most stable AUPRC — reliable baseline"],
    ["Spatial TransformerMIL", "0.859 ± 0.056", "0.465 ± 0.115",
     "Slight AUROC gain vs no-coords transformer; still poor and high-variance"],
    ["TransformerMIL\n(paper repro)", "0.850 ± 0.066", "0.508 ± 0.123",
     "Worst AUROC, highest variance — overfitting"],
]
ms_col_ws = [2.8, 2.2, 2.2, 4.75]
ms_col_ls = [0.35]
for w in ms_col_ws[:-1]:
    ms_col_ls.append(ms_col_ls[-1] + w + 0.05)

hdr_t = 1.3
row_h = 0.63
for j, (h, l, w) in enumerate(zip(ms_headers, ms_col_ls, ms_col_ws)):
    rect(sl, l, hdr_t, w, 0.45, fill=NAVY)
    txbox(sl, h, l + 0.05, hdr_t + 0.05, w - 0.1, 0.4,
          font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

row_bgs = [
    RGBColor(0xD4, 0xED, 0xDA),  # green — best
    RGBColor(0xE8, 0xF4, 0xF8),  # light blue
    RGBColor(0xEE, 0xF2, 0xFB),
    RGBColor(0xF0, 0xF0, 0xFB),
    RGBColor(0xF5, 0xF5, 0xF5),
    RGBColor(0xFD, 0xF0, 0xE8),  # orange-ish — spatial hurt
    RGBColor(0xF5, 0xF5, 0xF5),
    RGBColor(0xF9, 0xF1, 0xE7),  # warm neutral — slight transformer improvement, still weak
    RGBColor(0xFD, 0xE8, 0xE8),  # red — worst
]
for i, (row, rc) in enumerate(zip(ms_rows, row_bgs)):
    t = hdr_t + 0.45 + i * row_h
    for j, (cell, l, w) in enumerate(zip(row, ms_col_ls, ms_col_ws)):
        rect(sl, l, t, w, row_h - 0.04,
             fill=rc, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.4))
        is_metric = j in (1, 2)
        best = i == 0
        worst = i == 7
        col = (GREEN if best and is_metric else
               RED if worst and is_metric else DGRAY)
        txbox(sl, cell, l + 0.06, t + 0.04, w - 0.12, row_h - 0.12,
              font_size=12, bold=(is_metric and (best or worst)),
              color=col,
              align=PP_ALIGN.CENTER if is_metric else PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 23 — Hybrid attention deep-dive
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "The Winner: HybridAttentionMIL",
           "Best overall multisplit result: AUPRC 0.591 ± 0.058, AUROC 0.903 ± 0.035")

txbox(sl, "Architecture", 0.35, 1.35, 5.9, 0.38,
      font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "This is the direct multi-head extension of AttentionMIL: 2 learned attention heads instead of 1",
    "Per-head softmax → weighted sum → 1024-dim context vector per head",
    "Mean pool of full bag → 1024-dim global summary",
    "Concatenate: [head1 | head2 | mean] → 3×1024 = 3072-dim",
    "Params: 918,403 total vs 393,986 for single-head AttentionMIL",
    "Most of the increase is explicit: the classifier input grows from 1024 to 3072, "
     "adding 524,288 weights before the final head",
], l=0.35, t=1.82, w=5.9, h=3.6, font_size=14, indent_char="▸ ")

txbox(sl, "Open design questions", 0.35, 5.1, 5.9, 0.38,
      font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "How many heads are enough for CRC morphology: 2 is a strong start, but the clinical / embedding question is whether the task wants ~2, ~4, or more distinct tissue subtypes",
    "Exploration route: estimate natural embedding clusters, then ask whether heads align with recognisable tissue / morphology groups",
    "Training route: test diversity penalties, diversity-aware initialisation, and whether heads collapse without explicit pressure to separate",
    "Once head count is fixed, bag-level augmentation may matter more because it could improve coverage of rare morphologies seen by each head",
], l=0.35, t=5.55, w=5.9, h=1.55, font_size=14, indent_char="▸ ")

txbox(sl, "Multisplit results", 6.6, 1.35, 6.38, 0.38,
      font_size=16, bold=True, color=GREEN)
bullet_list(sl, [
    ("AUPRC:", "0.591 ± 0.058 — best overall and the metric to emphasise under class imbalance"),
    ("AUROC:", "0.903 ± 0.035 — also best overall, so the winner does not depend on the metric"),
    ("AUPRC variance:", "±0.058 vs ±0.128 for plain AttentionMIL — the hybrid is not just better on average, it is markedly more stable on the positive-class metric"),
    ("vs. AttentionMIL:", "AUPRC +0.059, AUROC +0.003, with much tighter AUPRC spread after adding a second head plus the mean branch"),
    ("Interpretation:", "Multiple heads are helping, but the next question is whether they are learning genuinely different morphologies or just a better optimization path"),
], l=6.6, t=1.82, w=6.38, h=5.5, font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 24 — Spatial coordinates
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Spatial Coordinate Encoding — A Mixed Story",
           "MLP coord encoder: why it was the right first choice, and what the results tell us")

# Three-column layout: MLP rationale | Results | What next
txbox(sl, "Why an MLP coordinate encoder?", 0.35, 1.28, 4.0, 0.38,
      font_size=14, bold=True, color=NAVY)
bullet_list(sl, [
    "No structural prior: unlike sinusoidal or RoPE encodings (which assume periodicity or "
     "rotation invariance), the MLP makes no assumption about which spatial patterns matter — "
     "it learns them from the task",
    "Continuous representation: no quantization artifacts from grid bins",
    "End-to-end: jointly optimized with aggregator — unused spatial signal "
     "gets suppressed by gradient descent",
    "Small footprint: 32-dim output is ~3% of the 1024-dim patch vector — "
     "a hint, not a dominant feature",
    "Honest weakness: needs labelled slides to learn meaningful spatial functions; "
     "at this dataset scale the gradient pressure on the coordinate branch is small — "
     "structured encodings (sinusoidal, RoPE) provide a useful prior without requiring learning",
], l=0.35, t=1.72, w=4.0, h=5.55, font_size=12, indent_char="▸ ")

rect(sl, 4.55, 1.28, 0.04, 5.9, fill=RGBColor(0xCC, 0xCC, 0xCC))  # divider

txbox(sl, "Multisplit results", 4.75, 1.28, 4.15, 0.38,
      font_size=14, bold=True, color=NAVY)
bullet_list(sl, [
    ("Attention + coords:", "AUROC 0.882 ± 0.041 — worse than plain (0.900 ± 0.032). "
     "AUPRC variance explodes. Spatial signal conflicts with position-agnostic "
     "attention scoring; the head cannot easily decouple location from content."),
    ("Hybrid + coords:", "AUROC 0.897 ± 0.040 vs 0.903 ± 0.035 non-spatial — "
     "essentially identical. Mean pool component absorbs the spatial noise, "
     "explaining why the hybrid is more robust to the addition."),
    ("Transformer (no coords):", "AUROC 0.850 ± 0.066, AUPRC 0.508 ± 0.123 — "
     "worst AUROC and highest variance overall."),
    ("Spatial TransformerMIL:", "AUROC 0.859 ± 0.056, AUPRC 0.465 ± 0.115 — "
     "slightly better AUROC than no-coords, but still poor and high-variance. "
     "The MLP coordinate branch does not rescue transformer instability."),
], l=4.75, t=1.72, w=4.15, h=5.55, font_size=12)

rect(sl, 9.1, 1.28, 0.04, 5.9, fill=RGBColor(0xCC, 0xCC, 0xCC))  # divider

txbox(sl, "Alternative encodings to try", 9.3, 1.28, 3.7, 0.38,
      font_size=14, bold=True, color=NAVY)
bullet_list(sl, [
    ("Sinusoidal 2D:", "Fixed-frequency sin/cos encodings — no learning needed, "
     "tested for years in NLP/ViT; good baseline for spatial periodicity"),
    ("Learned grid:", "Discretize slide into an N×N grid, embed cell index — "
     "simple but quantizes position"),
    ("RoPE:", "Rotary position embeddings applied to (x, y) — encodes relative "
     "rather than absolute position; may generalize better across slide sizes"),
    ("Region pooling:", "Tile the slide into macro-regions, pool within each, "
     "then aggregate across regions — intermediate granularity between "
     "patch-level and slide-level"),
    ("Key question:", "All alternatives trade learned flexibility for structured prior. "
     "The MLP result suggests the data cannot teach the spatial function reliably — "
     "a stronger prior is warranted."),
], l=9.3, t=1.72, w=3.7, h=5.55, font_size=12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 25 — Multisplit cohort stratification
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Cohort-Stratified Analysis",
           "4-model view: cohorts shown side by side for each model")

cohort_plot = make_cohort_grouped_plot(MULTI / "cohort_grouped_4models.png")
if cohort_plot and Path(cohort_plot).exists():
    img(sl, cohort_plot, 0.35, 1.3, 8.2)
    caption(sl, "Bars: mean cohort AUROC across 9 runs. Dots: split-level means.",
            0.35, 6.82, 8.2)

txbox(sl, "Takeaways", 8.85, 1.35, 4.1, 0.35, font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "The combined training scheme is fairly stable across cohorts for all 4 main models",
    "That is a real improvement over the earlier leaked setting, where cohort behaviour was much less trustworthy",
    "HybridAttentionMIL remains strongest overall without collapsing on one cohort to win the other",
    "The remaining cohort gap looks like ordinary dataset difficulty, not obvious evidence of split leakage",
], l=8.85, t=1.8, w=4.1, h=3.2, font_size=14, indent_char="▸ ")

txbox(sl, "Why this matters", 8.85, 5.2, 4.1, 0.35, font_size=16, bold=True, color=NAVY)
bullet_list(sl, [
    "We are training one combined model, not cohort-specific models",
    "Stable side-by-side cohort performance suggests the label reconciliation and case-grouping fixes are doing their job",
    "The right next step is external generalisation, not retreating to cohort-specific tuning too early",
], l=8.85, t=5.62, w=4.1, h=1.55, font_size=13, indent_char="▸ ")


# ═══════════════════════════════════════════════════════════════════════════════
# ── SECTION 8 ─────────────────────────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

sl = slide()
section_divider(sl, 8, "Summary & Future Directions",
                "What we learned and where to go next")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 25 — Key findings
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Key Findings", "")

findings = [
    ("1", "UNI embeddings are highly discriminative",
     "A simple mean-pooled linear probe achieves AUROC 0.860 on the fair-comparison split "
     "and 0.877 (mean) across 9 multi-split runs — exceeding the paper's 0.827. "
     "Representation quality of the foundation model is the dominant factor.",
     TEAL),
    ("2", "Hybrid attention is the clear multi-split winner",
     "HybridAttentionMIL (mean + 2 heads) achieves AUROC 0.903 ± 0.035 and AUPRC 0.591 ± 0.058 "
     "— best across all 9 models. Combining mean pooling with learned attention heads "
     "improves both performance and stability.",
     GREEN),
    ("3", "Complexity does not consistently win",
     "TransformerMIL underperforms (AUROC 0.850 ± 0.066), and adding MLP spatial coordinates "
     "only nudges it to 0.859 ± 0.056 while leaving variance high. Spatial coordinates hurt "
     "attention but are neutral for the hybrid. "
     "Data quality (leakage, label reconciliation) had larger impact than model choice.",
     NAVY),
    ("4", "Attention adds interpretability at manageable cost",
     "Attention maps reveal plausible tissue phenotype preferences. "
     "Multi-head attention separates concerns — each head can specialise. "
     "Feature-diverse sampling and mean-pool anchoring both reduce seed sensitivity.",
     ACCENT),
]
for i, (num, title, body, col) in enumerate(findings):
    t = 1.3 + i * 1.45
    rect(sl, 0.35, t, 0.6, 1.2, fill=col)
    txbox(sl, num, 0.35, t, 0.6, 1.2,
          font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, title, 1.1, t, 11.85, 0.42,
          font_size=17, bold=True, color=col)
    txbox(sl, body, 1.1, t + 0.43, 11.85, 0.9,
          font_size=14, color=DGRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 26 — Future directions
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Future Directions", "Immediate priorities and longer-horizon research")

# Left column — short term
rect(sl, 0.35, 1.3, 5.9, 5.7, fill=RGBColor(0xE8, 0xF4, 0xF8),
     line_color=TEAL, line_width=Pt(1.2))
txbox(sl, "Near-Term (next experiments)", 0.5, 1.43, 5.6, 0.42,
      font_size=16, bold=True, color=TEAL)
bullet_list(sl, [
    "HybridAttentionMIL validated ✓ — next: head count ablation (1/2/4/6), "
     "diversity penalty tuning, attention collapse diagnostics",
    "Spatial TransformerMIL did not rescue instability with an MLP coord encoder "
     "(AUROC 0.859 ± 0.056) — next spatial tests should use stronger priors",
    "Alternative spatial encodings: sinusoidal 2D positional embeddings, "
     "learned grid embeddings, RoPE-style rotary encodings applied to patch coords — "
     "the MLP coordinate encoder is a simple baseline, not the final answer",
    "Spatial context window: region-pooling (e.g. grid tiles) before attention — "
     "intermediate between patch-level and slide-level aggregation",
    "Bag-level augmentation: patch dropout, embedding perturbation, "
     "stain-normalisation proxy",
], l=0.5, t=1.93, w=5.6, h=4.7, font_size=14, indent_char="▸ ")

# Right column — longer horizon
rect(sl, 6.6, 1.3, 6.38, 5.7, fill=RGBColor(0xE8, 0xF8, 0xEE),
     line_color=GREEN, line_width=Pt(1.2))
txbox(sl, "Longer Horizon", 6.75, 1.43, 6.1, 0.42,
      font_size=16, bold=True, color=GREEN)
bullet_list(sl, [
    "Monitor cohort-stratified metrics to detect systematic gaps — "
     "per-cohort intervention (separate thresholds, separate heads) only if "
     "stratified multisplit plots show a consistent, meaningful divergence. "
     "Note: clinical deployment will not be SR1482 or SR386 — a third-cohort "
     "generalization test is the real bottleneck.",
    "False positive morphology audit: attention entropy as a reliability signal; "
     "flag high-confidence FP cases for pathologist review",
    "Additional tasks: KRAS/NRAS/BRAF mutation prediction, survival, "
     "multi-task learning",
    "Clinical interpretability: tissue phenotype cluster discovery via "
     "multi-head attention (expected clusters: tumour epi, TILs, stroma, mucin, necrosis)",
    "UNI fine-tuning if labelled case count increases substantially",
], l=6.75, t=1.93, w=6.1, h=4.7, font_size=14, indent_char="▸ ")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 27 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl, NAVY)
rect(sl, 0, 5.8, 13.33, 1.7, fill=TEAL)

txbox(sl, "Conclusion", 1.0, 1.0, 11.33, 0.9,
      font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider_line(sl, 2.1, l=3.0, w=7.33)

conclusions = [
    "Frozen UNI embeddings make MSI/MMR prediction tractable with simple aggregation.",
    "Mean pooling is a surprisingly strong baseline — AUROC 0.860 (fixed split), 0.877 (multi-split mean).",
    "HybridAttentionMIL (mean + 2 heads) is the multi-split winner: AUROC 0.903, AUPRC 0.591.",
    "Spatial coordinates: neutral for hybrid, harmful for plain attention, and only a slight AUROC gain for transformer without fixing instability.",
    "Data quality (case leakage, label reconciliation) had larger impact than model complexity.",
    "The framework is modular, reproducible, and production-ready for inference.",
]
for i, c in enumerate(conclusions):
    txbox(sl, f"▸  {c}", 1.0, 2.25 + i * 0.68, 11.33, 0.62,
          font_size=16, color=LIGHT, align=PP_ALIGN.CENTER)

txbox(sl, "Thank you", 1.0, 6.55, 11.33, 0.65,
      font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 28 — Backup: Calibration
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Backup: Calibration", "Temperature scaling applied post-hoc on validation set")

cal_path = OUT / "calibration.png"
if cal_path.exists():
    img(sl, cal_path, 0.5, 1.3, 12.0)
    caption(sl, "Reliability diagrams after temperature scaling · all three models",
            0.5, 6.8, 12.0)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 29 — Backup: Error distributions
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Backup: Error Distributions", "Score distributions by predicted outcome")

ed_path = OUT / "error_distributions.png"
if ed_path.exists():
    img(sl, ed_path, 0.5, 1.3, 12.0)
    caption(sl, "Predicted probability distributions by true label and model",
            0.5, 6.8, 12.0)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 30 — Backup: Appendix confusion matrices C & D
# ─────────────────────────────────────────────────────────────────────────────
sl = slide()
bg(sl)
header_bar(sl, "Backup: Appendix C & D Confusion Matrices")

cc = FIG / "appendix_c_confusion_matrices.png"
dc = FIG / "appendix_d_confusion_matrices.png"
if cc.exists():
    img(sl, cc, 0.35, 1.3, 6.1)
    caption(sl, "Appendix C: full-bag vs top-k=16 confusion matrices", 0.35, 6.82, 6.1)
if dc.exists():
    img(sl, dc, 6.85, 1.3, 6.1)
    caption(sl, "Appendix D: sampler ablation confusion matrices", 6.85, 6.82, 6.1)

# ─────────────────────────────────────────────────────────────────────────────
out_path = REPO / "surgen_mil_presentation.pptx"
prs.save(str(out_path))
print(f"\nSaved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
